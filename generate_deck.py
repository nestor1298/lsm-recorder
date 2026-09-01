#!/usr/bin/env python3
"""Generate SignaLab landing page presentation deck (.pptx)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Design Tokens ──────────────────────────────────────────────
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_DARK = RGBColor(0x31, 0x30, 0x7A)
INDIGO_LIGHT = RGBColor(0xEE, 0xF2, 0xFF)
CHARCOAL = RGBColor(0x11, 0x18, 0x27)
GRAY_600 = RGBColor(0x4B, 0x55, 0x63)
GRAY_400 = RGBColor(0x9C, 0xA3, 0xAF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORAL = RGBColor(0xF9, 0x73, 0x16)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ── Helpers ────────────────────────────────────────────────────
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT, spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Arial"
        p.alignment = align
        p.space_after = Pt(font_size * (spacing - 1))
    return txBox


def add_slide_number(slide, num, total=8):
    add_text(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.6),
             Inches(1), Inches(0.4), f"{num} / {total}",
             font_size=11, color=GRAY_400, align=PP_ALIGN.RIGHT)


def add_avatar_placeholder(slide, left, top, size, label="Signa"):
    """Draws a circle with 'Signa' label as avatar placeholder."""
    circle = add_circle(slide, left, top, size, INDIGO_LIGHT)
    # Add a smaller indigo circle inside
    inner_offset = int(size * 0.15)
    inner_size = int(size * 0.7)
    inner = add_circle(slide, left + inner_offset, top + inner_offset, inner_size, INDIGO)
    # Label
    add_text(slide, left, top + size + Inches(0.1), size, Inches(0.3),
             label, font_size=10, color=GRAY_400, align=PP_ALIGN.CENTER)
    return circle


# ── SLIDE 1: Hero ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)

# Indigo accent bar at top
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), INDIGO)

# Left content
add_text(slide, Inches(1.2), Inches(1.5), Inches(4), Inches(0.6),
         "SignaLab", font_size=20, color=INDIGO, bold=True)

add_text(slide, Inches(1.2), Inches(2.2), Inches(6.5), Inches(1.5),
         "Documenting Mexican Sign\nLanguage, One Hand at a Time",
         font_size=44, color=CHARCOAL, bold=True)

add_text(slide, Inches(1.2), Inches(4.2), Inches(6), Inches(1),
         "Computational linguistics research advancing linguistic rights\nin education and the future of work for the Deaf community.",
         font_size=18, color=GRAY_600)

# CTA button
btn = add_rounded_rect(slide, Inches(1.2), Inches(5.5), Inches(3.2), Inches(0.7), INDIGO)
tf = btn.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Start Contributing  →"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Arial"
p.alignment = PP_ALIGN.CENTER
tf.paragraphs[0].space_before = Pt(6)

# Right side — avatar placeholder
add_avatar_placeholder(slide, Inches(9), Inches(1.8), Inches(3.5), "[ Signa — waving hello ]")

# URL
add_text(slide, Inches(1.2), Inches(6.5), Inches(5), Inches(0.4),
         "lsm-recorder.vercel.app", font_size=13, color=GRAY_400)

add_slide_number(slide, 1)


# ── SLIDE 2: Mission ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, INDIGO)

add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, INDIGO)

add_text(slide, Inches(1.5), Inches(1.2), Inches(4), Inches(0.5),
         "OUR MISSION", font_size=14, color=INDIGO_LIGHT, bold=True)

add_text(slide, Inches(1.5), Inches(2), Inches(10), Inches(2.5),
         "Build the first open-source phonological\ndocumentation system for Mexican\nSign Language (LSM)",
         font_size=40, color=WHITE, bold=True)

add_multiline(slide, Inches(1.5), Inches(4.8), Inches(8), Inches(2), [
    "We combine the Cruz Aldrete (2008) phonological model with modern",
    "computer vision to create a complete pipeline: from video recording,",
    "through linguistic annotation, to 3D avatar animation.",
    "",
    "Every sign. Every handshape. Every movement — documented and preserved."
], font_size=17, color=INDIGO_LIGHT, spacing=1.4)

add_slide_number(slide, 2)


# ── SLIDE 3: Features ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text(slide, Inches(1.2), Inches(0.6), Inches(6), Inches(0.5),
         "WHAT SIGNALAB DOES", font_size=14, color=INDIGO, bold=True)

add_text(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(0.8),
         "Four tools, one platform", font_size=36, color=CHARCOAL, bold=True)

features = [
    ("Record", "Capture sign language videos\nwith guided prompts for each\nof the 101 handshape\nconfigurations.", INDIGO, "📹"),
    ("Catalog", "Browse the complete Cruz\nAldrete inventory — 101 CMs\norganized by frequency tier\nand finger group.", BLUE, "📖"),
    ("Annotate", "Mark phonological segments\nusing the PSHR timeline:\nPreparation, Stroke, Hold,\nRetraction.", GREEN, "✏️"),
    ("Body Map", "Select articulation locations\non an interactive body\nsilhouette — 80 UB points\nacross 7 body regions.", CORAL, "🧍"),
]

card_w = Inches(2.6)
card_h = Inches(3.6)
start_x = Inches(1.2)
gap = Inches(0.35)

for i, (title, desc, color, emoji) in enumerate(features):
    x = start_x + i * (card_w + gap)
    y = Inches(2.3)

    # Card background
    card = add_rounded_rect(slide, x, y, card_w, card_h, WHITE)
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    card.line.width = Pt(1)

    # Color accent bar
    add_rect(slide, x, y, card_w, Inches(0.06), color)

    # Emoji placeholder
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.6),
             emoji, font_size=28, align=PP_ALIGN.LEFT)

    # Title
    add_text(slide, x + Inches(0.3), y + Inches(0.9), Inches(2), Inches(0.4),
             title, font_size=20, color=CHARCOAL, bold=True)

    # Description
    add_text(slide, x + Inches(0.3), y + Inches(1.4), Inches(2.1), Inches(2),
             desc, font_size=13, color=GRAY_600)

add_slide_number(slide, 3)


# ── SLIDE 4: For the Community ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Left half — lavender background
add_rect(slide, 0, 0, Inches(6.666), SLIDE_H, INDIGO_LIGHT)

add_text(slide, Inches(1.2), Inches(1), Inches(4), Inches(0.5),
         "FOR EDUCATORS & THE DEAF COMMUNITY", font_size=12, color=INDIGO, bold=True)

add_text(slide, Inches(1.2), Inches(1.7), Inches(5), Inches(1.2),
         "Visual-first tools\nbuilt for you",
         font_size=36, color=CHARCOAL, bold=True)

bullets_community = [
    "→  Record signs at your own pace with guided prompts",
    "→  Every tool is visual — no linguistic jargon required",
    "→  Contribute to the largest LSM video corpus",
    "→  Help preserve and standardize LSM for education",
    "→  Your recordings train AI to understand sign language",
]
add_multiline(slide, Inches(1.2), Inches(3.3), Inches(5), Inches(3),
              bullets_community, font_size=16, color=GRAY_600, spacing=1.6)

# Right half — avatar placeholder
add_avatar_placeholder(slide, Inches(8.5), Inches(2), Inches(3), "[ Signa — teaching ]")

add_text(slide, Inches(7.5), Inches(5.5), Inches(5), Inches(1),
         "\"Nothing about us without us.\"\nLanguage rights start with documentation.",
         font_size=15, color=GRAY_400, align=PP_ALIGN.CENTER)

add_slide_number(slide, 4)


# ── SLIDE 5: For Researchers ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text(slide, Inches(1.2), Inches(0.6), Inches(6), Inches(0.5),
         "FOR RESEARCHERS & LINGUISTS", font_size=12, color=INDIGO, bold=True)

add_text(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(0.8),
         "Phonological precision meets modern tooling",
         font_size=32, color=CHARCOAL, bold=True)

# Left column
add_text(slide, Inches(1.2), Inches(2.4), Inches(3.5), Inches(0.4),
         "5-Channel Notation", font_size=18, color=INDIGO, bold=True)

notation_items = [
    "CM — Configuración Manual (handshape)",
    "UB — Ubicación (body location)",
    "MV — Movimiento (movement path & type)",
    "OR — Orientación (palm & finger direction)",
    "RNM — Rasgos No Manuales (face & head)",
]
add_multiline(slide, Inches(1.2), Inches(3), Inches(4.5), Inches(2.5),
              notation_items, font_size=14, color=GRAY_600, spacing=1.5)

# Right column
add_text(slide, Inches(7), Inches(2.4), Inches(4), Inches(0.4),
         "Data & Export", font_size=18, color=INDIGO, bold=True)

data_items = [
    "PSHR timeline segmentation",
    "JSON export for every annotation",
    "Cruz Aldrete (2008) model compliance",
    "MediaPipe hand keypoint extraction",
    "Open-source — all data is yours",
]
add_multiline(slide, Inches(7), Inches(3), Inches(5), Inches(2.5),
              data_items, font_size=14, color=GRAY_600, spacing=1.5)

# Notation example box
box = add_rounded_rect(slide, Inches(1.2), Inches(5.4), Inches(11), Inches(1.2), RGBColor(0xF9, 0xFA, 0xFB))
box.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
box.line.width = Pt(1)

add_text(slide, Inches(1.5), Inches(5.55), Inches(10.5), Inches(0.9),
         '{ "cm_id": 1,  "phase": "STROKE",  "location": "Frente",  "contour": "ARC",  "palm": "FORWARD" }',
         font_size=14, color=INDIGO_DARK, font_name="Courier New")

add_slide_number(slide, 5)


# ── SLIDE 6: Impact Numbers ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, INDIGO)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, INDIGO)

add_text(slide, Inches(1.2), Inches(1), Inches(10), Inches(0.5),
         "THE NUMBERS", font_size=14, color=INDIGO_LIGHT, bold=True)

add_text(slide, Inches(1.2), Inches(1.6), Inches(10), Inches(0.8),
         "Building at scale from day one",
         font_size=36, color=WHITE, bold=True)

stats = [
    ("101", "Handshape\nConfigurations", "Complete Cruz Aldrete\ninventory digitized"),
    ("80", "Body\nLocations", "Interactive articulation\npoints mapped to SVG"),
    ("5", "Annotation\nChannels", "CM, UB, MV, OR, RNM\nper segment"),
    ("4", "Frequency\nTiers", "Prioritized recording\nby usage frequency"),
]

stat_w = Inches(2.5)
stat_start = Inches(1.2)
stat_gap = Inches(0.5)

for i, (number, label, detail) in enumerate(stats):
    x = stat_start + i * (stat_w + stat_gap)
    y = Inches(3)

    # Number
    add_text(slide, x, y, stat_w, Inches(1),
             number, font_size=64, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

    # Label
    add_text(slide, x, y + Inches(1.2), stat_w, Inches(0.8),
             label, font_size=18, color=WHITE, bold=True)

    # Detail
    add_text(slide, x, y + Inches(2.1), stat_w, Inches(0.8),
             detail, font_size=12, color=INDIGO_LIGHT)

add_slide_number(slide, 6)


# ── SLIDE 7: The Technology / Pipeline ─────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text(slide, Inches(1.2), Inches(0.6), Inches(6), Inches(0.5),
         "THE PIPELINE", font_size=14, color=INDIGO, bold=True)

add_text(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(0.8),
         "From video to avatar — end to end",
         font_size=36, color=CHARCOAL, bold=True)

# Pipeline steps
steps = [
    ("1", "Record", "Capture video of\neach handshape\nconfiguration", INDIGO),
    ("2", "Annotate", "Segment & label\nusing LSM-PN\n5-channel notation", GREEN),
    ("3", "Extract", "MediaPipe hand\nkeypoints + body\npose estimation", BLUE),
    ("4", "Classify", "Train ML models\non phonological\nfeatures", AMBER),
    ("5", "Animate", "Drive 3D avatar\nhands from\nnotation data", CORAL),
]

step_w = Inches(2)
step_start = Inches(1)
step_gap = Inches(0.3)
arrow_w = Inches(0.3)

for i, (num, title, desc, color) in enumerate(steps):
    x = step_start + i * (step_w + step_gap)
    y = Inches(2.6)

    # Step circle
    circle = add_circle(slide, x + Inches(0.65), y, Inches(0.7), color)
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)

    # Title
    add_text(slide, x, y + Inches(1), step_w, Inches(0.4),
             title, font_size=18, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER)

    # Desc
    add_text(slide, x, y + Inches(1.5), step_w, Inches(1.2),
             desc, font_size=12, color=GRAY_600, align=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < len(steps) - 1:
        add_text(slide, x + step_w, y + Inches(0.15), Inches(0.4), Inches(0.5),
                 "→", font_size=24, color=GRAY_400, align=PP_ALIGN.CENTER)

# Bottom note
box = add_rounded_rect(slide, Inches(1.2), Inches(5.6), Inches(11), Inches(1), INDIGO_LIGHT)
box.line.fill.background()

add_text(slide, Inches(1.5), Inches(5.75), Inches(10.5), Inches(0.7),
         "Phase 1 (current): Record + Annotate web app    |    Phase 2: Extract + Classify ML pipeline    |    Phase 3: Avatar animation",
         font_size=14, color=INDIGO_DARK, align=PP_ALIGN.CENTER)

add_slide_number(slide, 7)


# ── SLIDE 8: CTA / Closing ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, INDIGO)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, INDIGO)

# Avatar placeholder
add_avatar_placeholder(slide, Inches(5.6), Inches(0.8), Inches(2.2), "[ Signa — celebrating ]")

add_text(slide, Inches(2), Inches(3.3), Inches(9.5), Inches(1.2),
         "Ready to help document\nMexican Sign Language?",
         font_size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# CTA button
btn = add_rounded_rect(slide, Inches(4.5), Inches(5), Inches(4.5), Inches(0.8), WHITE)
tf = btn.text_frame
p = tf.paragraphs[0]
p.text = "Open SignaLab  →"
p.font.size = Pt(22)
p.font.color.rgb = INDIGO
p.font.bold = True
p.font.name = "Arial"
p.alignment = PP_ALIGN.CENTER
tf.paragraphs[0].space_before = Pt(8)

add_text(slide, Inches(3), Inches(6.1), Inches(7.5), Inches(0.4),
         "lsm-recorder.vercel.app", font_size=16, color=INDIGO_LIGHT, align=PP_ALIGN.CENTER)

add_text(slide, Inches(3), Inches(6.6), Inches(7.5), Inches(0.4),
         "SignaLab  ·  Computational Linguistics for Linguistic Rights",
         font_size=12, color=INDIGO_LIGHT, align=PP_ALIGN.CENTER)

add_slide_number(slide, 8)


# ── Save ───────────────────────────────────────────────────────
output = "/Users/nestoraldreteochoa/Downloads/lsm-recorder/SignaLab_Landing.pptx"
prs.save(output)
print(f"✓ Saved {output}")
print(f"  {len(prs.slides)} slides, 16:9 widescreen")
